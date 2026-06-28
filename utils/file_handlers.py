"""
utils/file_handlers.py - File Text Extraction Engine
AI Study Notes Generator
-----------------------------------------
Handles text extraction from:
- PDF files (via pdfplumber + PyPDF2 fallback)
- DOCX files (via python-docx)
- PPTX files (via python-pptx)
- TXT files (plain read)
"""

import os
import logging

# PDF Libraries
# pyrefly: ignore [missing-import]
import pdfplumber
# pyrefly: ignore [missing-import]
import PyPDF2

# Document Libraries
# pyrefly: ignore [missing-import]
from docx import Document

# Presentation Libraries
# pyrefly: ignore [missing-import]
from pptx import Presentation

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


# ─────────────────────────────────────────
# PDF Text Extraction
# ─────────────────────────────────────────

def extract_text_from_pdf(file_path: str) -> str:
    """
    Extract text from a PDF file.
    Uses pdfplumber first (better accuracy),
    falls back to PyPDF2 if pdfplumber fails.

    Args:
        file_path: Path to the PDF file

    Returns:
        Extracted text as a string
    """
    text = ""

    # ── Method 1: pdfplumber (preferred) ──
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

        if text.strip():
            logger.info(f"✅ PDF extracted via pdfplumber: {len(text)} chars")
            return text.strip()

    except Exception as e:
        logger.warning(f"⚠️ pdfplumber failed: {e}. Trying PyPDF2...")

    # ── Method 2: PyPDF2 (fallback) ──
    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

        if text.strip():
            logger.info(f"✅ PDF extracted via PyPDF2: {len(text)} chars")
            return text.strip()

    except Exception as e:
        logger.error(f"❌ PyPDF2 also failed: {e}")

    raise ValueError("Could not extract text from PDF. The file may be scanned or image-based.")


# ─────────────────────────────────────────
# DOCX Text Extraction
# ─────────────────────────────────────────

def extract_text_from_docx(file_path: str) -> str:
    """
    Extract text from a DOCX Word document.
    Extracts paragraphs, tables, and headers.

    Args:
        file_path: Path to the DOCX file

    Returns:
        Extracted text as a string
    """
    try:
        doc = Document(file_path)
        text_parts = []

        # ── Extract Paragraphs ──
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text.strip())

        # ── Extract Tables ──
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip()
                    for cell in row.cells
                    if cell.text.strip()
                )
                if row_text:
                    text_parts.append(row_text)

        text = "\n".join(text_parts)

        if not text.strip():
            raise ValueError("No text found in DOCX file.")

        logger.info(f"✅ DOCX extracted: {len(text)} chars")
        return text.strip()

    except Exception as e:
        logger.error(f"❌ DOCX extraction failed: {e}")
        raise ValueError(f"Could not extract text from DOCX: {str(e)}")


# ─────────────────────────────────────────
# PPTX Text Extraction
# ─────────────────────────────────────────

def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract text from a PPTX PowerPoint file.
    Extracts text from all slides, shapes, and notes.

    Args:
        file_path: Path to the PPTX file

    Returns:
        Extracted text as a string
    """
    try:
        prs = Presentation(file_path)
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []

            # ── Extract text from all shapes ──
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            # ── Extract speaker notes ──
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    slide_texts.append(f"[Notes]: {notes.strip()}")

            if slide_texts:
                text_parts.append(f"--- Slide {slide_num} ---")
                text_parts.extend(slide_texts)

        text = "\n".join(text_parts)

        if not text.strip():
            raise ValueError("No text found in PPTX file.")

        logger.info(f"✅ PPTX extracted: {len(text)} chars from {len(prs.slides)} slides")
        return text.strip()

    except Exception as e:
        logger.error(f"❌ PPTX extraction failed: {e}")
        raise ValueError(f"Could not extract text from PPTX: {str(e)}")


# ─────────────────────────────────────────
# TXT Text Extraction
# ─────────────────────────────────────────

def extract_text_from_txt(file_path: str) -> str:
    """
    Extract text from a plain TXT file.
    Tries UTF-8 first, falls back to latin-1.

    Args:
        file_path: Path to the TXT file

    Returns:
        Extracted text as a string
    """
    # ── Try UTF-8 first ──
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()

        if text.strip():
            logger.info(f"✅ TXT extracted (UTF-8): {len(text)} chars")
            return text.strip()

    except UnicodeDecodeError:
        pass

    # ── Fallback to latin-1 ──
    try:
        with open(file_path, "r", encoding="latin-1") as f:
            text = f.read()

        if text.strip():
            logger.info(f"✅ TXT extracted (latin-1): {len(text)} chars")
            return text.strip()

    except Exception as e:
        logger.error(f"❌ TXT extraction failed: {e}")
        raise ValueError(f"Could not read TXT file: {str(e)}")

    raise ValueError("TXT file appears to be empty.")


# ─────────────────────────────────────────
# Main Router Function
# ─────────────────────────────────────────

def extract_text(file_path: str) -> str:
    """
    Main router — detects file type and calls
    the appropriate extraction function.

    Args:
        file_path: Path to any supported file

    Returns:
        Extracted text as a string

    Raises:
        ValueError: If file type is unsupported or extraction fails
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    # Get file extension
    _, ext = os.path.splitext(file_path)
    ext = ext.lower().strip(".")

    # Route to correct extractor
    extractors = {
        "pdf":  extract_text_from_pdf,
        "docx": extract_text_from_docx,
        "pptx": extract_text_from_pptx,
        "txt":  extract_text_from_txt,
    }

    if ext not in extractors:
        raise ValueError(
            f"Unsupported file type: .{ext}. "
            f"Supported types: {', '.join(extractors.keys())}"
        )

    logger.info(f"📄 Extracting text from {ext.upper()} file: {file_path}")
    return extractors[ext](file_path)


# ─────────────────────────────────────────
# File Validation
# ─────────────────────────────────────────

def validate_file(filename: str, file_size: int, max_size: int) -> dict:
    """
    Validate uploaded file before processing.

    Args:
        filename:  Original filename
        file_size: File size in bytes
        max_size:  Maximum allowed size in bytes

    Returns:
        dict with 'valid' bool and 'error' message
    """
    # Check filename
    if not filename or filename.strip() == "":
        return {"valid": False, "error": "No filename provided."}

    # Check extension
    _, ext = os.path.splitext(filename)
    ext = ext.lower().strip(".")
    allowed = ["pdf", "docx", "pptx", "txt"]

    if ext not in allowed:
        return {
            "valid": False,
            "error": f"File type '.{ext}' is not supported. "
                     f"Please upload: {', '.join(allowed)}"
        }

    # Check file size
    if file_size > max_size:
        max_mb = max_size / (1024 * 1024)
        return {
            "valid": False,
            "error": f"File is too large. Maximum size is {max_mb:.0f} MB."
        }

    return {"valid": True, "error": None}
